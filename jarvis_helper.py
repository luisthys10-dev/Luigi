#!/usr/bin/env python3
"""JARVIS Python helper - callable via: py C:\\Jarvis\\jarvis_helper.py [function] [args]"""
import sys
import json
import base64
import os

SCREENSHOT_PATH = r"C:\Jarvis\screenshot.png"


def make_word(path, content):
    from docx import Document
    doc = Document()
    for line in content.split("\\n"):
        doc.add_paragraph(line)
    doc.save(path)
    print(json.dumps({"success": True, "path": path}))


def make_excel(path, data):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for row_text in data.split("\\n"):
        cells = [c.strip() for c in row_text.split(",")]
        ws.append(cells)
    wb.save(path)
    print(json.dumps({"success": True, "path": path}))


def make_pptx(path, content):
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    lines = content.split("\\n")
    title_shape.text = lines[0] if lines else "JARVIS"
    body_shape.text = "\n".join(lines[1:]) if len(lines) > 1 else content
    prs.save(path)
    print(json.dumps({"success": True, "path": path}))


def take_screenshot():
    from PIL import ImageGrab
    img = ImageGrab.grab()
    os.makedirs(os.path.dirname(SCREENSHOT_PATH), exist_ok=True)
    img.save(SCREENSHOT_PATH)
    with open(SCREENSHOT_PATH, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    print(json.dumps({"success": True, "base64": b64, "path": SCREENSHOT_PATH}))


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Geen functie opgegeven"}))
        sys.exit(1)

    func = sys.argv[1]

    try:
        if func == "make_word":
            if len(sys.argv) < 4:
                raise ValueError("Gebruik: make_word <pad> <inhoud>")
            make_word(sys.argv[2], sys.argv[3])
        elif func == "make_excel":
            if len(sys.argv) < 4:
                raise ValueError("Gebruik: make_excel <pad> <data>")
            make_excel(sys.argv[2], sys.argv[3])
        elif func == "make_pptx":
            if len(sys.argv) < 4:
                raise ValueError("Gebruik: make_pptx <pad> <inhoud>")
            make_pptx(sys.argv[2], sys.argv[3])
        elif func == "take_screenshot":
            take_screenshot()
        else:
            print(json.dumps({"error": f"Onbekende functie: {func}"}))
            sys.exit(1)
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)
